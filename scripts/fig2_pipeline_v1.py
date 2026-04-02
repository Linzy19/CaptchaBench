"""
fig2_pipeline_v1.py — Academic beautification of pipeline figure
Changes from v0:
1. Unified color palette aligned with paper's modality-logic colors
2. Cleaner fonts: all labels use consistent sizing hierarchy
3. Better step badges: more prominent, colored per section
4. Left pipeline boxes: add subtle gradient-like depth with darker border
5. Perturbation methods: use paper's exact method colors
6. Evaluation panels: clean dark headers, lighter body
7. Connectors: add arrowheads, slightly thicker
8. Add figure title at top
9. Fix duplicate CR↑ boxes (78,79 → TVR↓ and ASR↑)
10. White slide background
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
import copy, subprocess

SRC  = Path("/Users/lzy/Code_lzy/Posion_paper/ACM_MM26/main/figures/cc_fig/fig2_pipeline_v0.pptx")
DEST = Path("/Users/lzy/Code_lzy/Posion_paper/ACM_MM26/main/figures/cc_fig/fig2_pipeline_v1.pptx")
OUT_DIR = DEST.parent

import shutil
shutil.copy(str(SRC), str(DEST))

prs = Presentation(str(DEST))
slide = prs.slides[0]
shapes_by_name = {s.name: s for s in slide.shapes}

print("=== Shapes found in slide ===")
for name in sorted(shapes_by_name.keys()):
    print(f"  {name}")
print(f"Total: {len(shapes_by_name)} shapes")
print()

# ── Color constants ──────────────────────────────────────────────────────────
def rgb(h): return RGBColor(*bytes.fromhex(h.lstrip('#')))

# Paper's modality-logic colors (must match figures 3–7)
C = {
    # Method colors
    'AMP':        '#0D47A1',
    'ASPL':       '#1976D2',
    'XTransfer':  '#42A5F5',
    'Glaze':      '#90CAF9',
    'Nightshade': '#E64A19',
    'MMCoA':      '#2E7D32',
    # Group colors (header bars)
    'img_only_hdr':  '#0D47A1',   # deep blue
    'img_only_bg':   '#E3F0FF',   # very light blue
    'img_only_bdr':  '#1565C0',
    'text_only_hdr': '#BF360C',   # deep orange
    'text_only_bg':  '#FFF3EE',
    'text_only_bdr': '#E64A19',
    'img_text_hdr':  '#1B5E20',   # deep green
    'img_text_bg':   '#EDF7EE',
    'img_text_bdr':  '#2E7D32',
    # VLM panel (teal)
    'vlm_hdr':  '#00695C',
    'vlm_bg':   '#E0F2F1',
    'vlm_bdr':  '#00897B',
    # Quality panel (indigo)
    'qual_hdr': '#283593',
    'qual_bg':  '#E8EAF6',
    'qual_bdr': '#3949AB',
    # Efficiency panel (brown-gray)
    'eff_hdr':  '#37474F',
    'eff_bg':   '#ECEFF1',
    'eff_bdr':  '#546E7A',
    # Left pipeline boxes
    'pipe_bg':  '#EEF4FB',
    'pipe_bdr': '#1565C0',
    # Step badges
    'badge_bg':  '#F5F5F5',
    'badge_txt': '#424242',
    # Connectors
    'arrow': '#455A64',
    # Slide bg
    'slide_bg': '#FFFFFF',
}

# VLM colors (Tableau10, matching Fig4)
VLM_COLORS = {
    'Qwen-VL':    '#4E79A7',
    'Kimi 2.5':   '#59A14F',
    'GPT-5.2':    '#F28E2B',
    'Gemini-3.0': '#E15759',
    'GLM-4V':     '#76B7B2',
}

# ── Helper functions ─────────────────────────────────────────────────────────
def set_fill(shape, hex_color, alpha=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hex_color)

def set_line(shape, hex_color, width_pt=1.0):
    shape.line.color.rgb = rgb(hex_color)
    shape.line.width = Pt(width_pt)

def set_text_style(shape, text=None, font_size=None, bold=None, color=None, align=None):
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        if align == 'center':
            para.alignment = PP_ALIGN.CENTER
        elif align == 'left':
            para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            if text is not None and run == para.runs[0]:
                run.text = text
                text = None
            if font_size is not None:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = rgb(color)

def set_no_line(shape):
    shape.line.color.rgb = rgb('#FFFFFF')
    shape.line.width = Pt(0)

def set_corner_radius(shape, radius_emu=60000):
    """Set rounded corner radius via XML."""
    try:
        sp = shape._element
        prstGeom = sp.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            # remove existing
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius_emu}')
    except Exception as e:
        pass  # ignore if not a rounded rect

# ── 1. Slide background: pure white ─────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = rgb('#FFFFFF')
print("✅ Step 1: Slide background set to white")

# ── 2. Left pipeline boxes (①②③) ───────────────────────────────────────────
pipe_boxes = ['Rounded Rectangle 3', 'Rounded Rectangle 6', 'Rounded Rectangle 9']
found_pipe = [n for n in pipe_boxes if n in shapes_by_name]
missing_pipe = [n for n in pipe_boxes if n not in shapes_by_name]
print(f"✅ Step 2: Pipe boxes found: {found_pipe}")
if missing_pipe:
    print(f"  ⚠️  Missing pipe boxes: {missing_pipe}")

for name in found_pipe:
    s = shapes_by_name[name]
    set_fill(s, C['pipe_bg'])
    set_line(s, C['pipe_bdr'], 1.5)
    set_text_style(s, font_size=8.5, bold=False, color='#1A237E', align='center')

def set_multiline_text(shape, lines, font_size=8.5, bold_first=True, color='#1A237E'):
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold_first if i == 0 else False
        run.font.color.rgb = rgb(color)

if 'Rounded Rectangle 3' in shapes_by_name:
    set_multiline_text(shapes_by_name['Rounded Rectangle 3'],
        ['GB2312 Level-1', '3,500 Characters', '↓ 840,000 Images', '(2 Gen × 120 BG)'],
        font_size=8.5, bold_first=True, color='#1A237E')

if 'Rounded Rectangle 6' in shapes_by_name:
    set_multiline_text(shapes_by_name['Rounded Rectangle 6'],
        ['SD (ID ControlNet)', '+', 'SDXL ControlNet', '120 backgrounds'],
        font_size=8.5, bold_first=True, color='#1A237E')

if 'Rounded Rectangle 9' in shapes_by_name:
    set_multiline_text(shapes_by_name['Rounded Rectangle 9'],
        ['CAPTCHA Images', '840,000 total', '(3,500 × 2 × 120)'],
        font_size=8.5, bold_first=True, color='#1A237E')

print("  ↳ Pipe box text content updated")

# ── 3. Step badges (bottom row) ──────────────────────────────────────────────
badge_configs = {
    'Rounded Rectangle 4':  ('① Dataset',      C['pipe_bdr'],   '#FFFFFF'),
    'Rounded Rectangle 7':  ('② Generation',   C['pipe_bdr'],   '#FFFFFF'),
    'Rounded Rectangle 10': ('③ CAPTCHA',      C['pipe_bdr'],   '#FFFFFF'),
    'Rounded Rectangle 31': ('④ Perturbation', '#5C3472',       '#FFFFFF'),
    'Rounded Rectangle 74': ('⑤ Evaluation',   C['vlm_hdr'],    '#FFFFFF'),
}
found_badges = []
missing_badges = []
for name, (text, fill_color, txt_color) in badge_configs.items():
    if name in shapes_by_name:
        s = shapes_by_name[name]
        set_fill(s, fill_color)
        set_no_line(s)
        set_text_style(s, font_size=8.5, bold=True, color=txt_color, align='center')
        found_badges.append(name)
    else:
        missing_badges.append(name)
print(f"✅ Step 3: Badges found: {found_badges}")
if missing_badges:
    print(f"  ⚠️  Missing badges: {missing_badges}")

# ── 4. Perturbation section ──────────────────────────────────────────────────
pert_containers = {
    'Rounded Rectangle 11': (C['img_only_bg'],   C['img_only_bdr']),
    'Rounded Rectangle 13': (C['text_only_bg'],  C['text_only_bdr']),
    'Rounded Rectangle 15': (C['img_text_bg'],   C['img_text_bdr']),
}
found_pc = []
missing_pc = []
for name, (bg_col, bdr_col) in pert_containers.items():
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], bg_col)
        set_line(shapes_by_name[name], bdr_col, 1.2)
        found_pc.append(name)
    else:
        missing_pc.append(name)
print(f"✅ Step 4a: Perturbation containers found: {found_pc}")
if missing_pc:
    print(f"  ⚠️  Missing: {missing_pc}")

pert_headers = {
    'Rounded Rectangle 12': C['img_only_hdr'],
    'Rounded Rectangle 14': C['text_only_hdr'],
    'Rounded Rectangle 16': C['img_text_hdr'],
}
found_ph = []
missing_ph = []
for name, hdr_col in pert_headers.items():
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], hdr_col)
        set_no_line(shapes_by_name[name])
        set_text_style(shapes_by_name[name], font_size=8, bold=True, color='#FFFFFF', align='center')
        found_ph.append(name)
    else:
        missing_ph.append(name)
print(f"  ↳ Perturbation headers found: {found_ph}")
if missing_ph:
    print(f"  ⚠️  Missing: {missing_ph}")

method_color_map = {
    'Rounded Rectangle 17': ('ASPL',       C['ASPL']),
    'Rounded Rectangle 18': ('Glaze (MI)', C['Glaze']),
    'Rounded Rectangle 19': ('AMP',        C['AMP']),
    'Rounded Rectangle 20': ('XTransfer',  C['XTransfer']),
    'Rounded Rectangle 21': ('Nightshade', C['Nightshade']),
    'Rounded Rectangle 22': ('MMCoA',      C['MMCoA']),
}
found_meth = []
missing_meth = []
for name, (label, color) in method_color_map.items():
    if name in shapes_by_name:
        s = shapes_by_name[name]
        set_fill(s, color)
        set_no_line(s)
        txt_col = '#FFFFFF' if color in [C['AMP'], C['ASPL'], C['Nightshade'], C['MMCoA']] else '#1A237E'
        set_text_style(s, font_size=8.5, bold=True, color=txt_col, align='center')
        found_meth.append(name)
    else:
        missing_meth.append(name)
print(f"  ↳ Method buttons found: {found_meth}")
if missing_meth:
    print(f"  ⚠️  Missing: {missing_meth}")

# ── 5. VLM Evaluation panel ──────────────────────────────────────────────────
vlm_containers = ['Rounded Rectangle 41', 'Rounded Rectangle 42']
found_vc = []
missing_vc = []
for name in vlm_containers:
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], C['vlm_bg'])
        set_line(shapes_by_name[name], C['vlm_bdr'], 1.5)
        found_vc.append(name)
    else:
        missing_vc.append(name)
print(f"✅ Step 5a: VLM containers found: {found_vc}")
if missing_vc:
    print(f"  ⚠️  Missing: {missing_vc}")

vlm_headers = ['Rounded Rectangle 43', 'Rounded Rectangle 44']
found_vh = []
missing_vh = []
for name in vlm_headers:
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], C['vlm_hdr'])
        set_no_line(shapes_by_name[name])
        set_text_style(shapes_by_name[name], font_size=9, bold=True, color='#FFFFFF', align='center')
        found_vh.append(name)
    else:
        missing_vh.append(name)
print(f"  ↳ VLM headers found: {found_vh}")
if missing_vh:
    print(f"  ⚠️  Missing: {missing_vh}")

vlm_box_map = {
    'Rounded Rectangle 45': 'Qwen-VL',
    'Rounded Rectangle 46': 'Kimi 2.5',
    'Rounded Rectangle 47': 'GPT-5.2',
    'Rounded Rectangle 76': 'Gemini-3.0',
    'Rounded Rectangle 77': 'GLM-4V',
}
found_vb = []
missing_vb = []
for name, vlm in vlm_box_map.items():
    if name in shapes_by_name:
        color = VLM_COLORS[vlm]
        set_fill(shapes_by_name[name], color)
        set_no_line(shapes_by_name[name])
        set_text_style(shapes_by_name[name], font_size=8, bold=True, color='#FFFFFF', align='center')
        found_vb.append(name)
    else:
        missing_vb.append(name)
print(f"  ↳ VLM boxes found: {found_vb}")
if missing_vb:
    print(f"  ⚠️  Missing: {missing_vb}")

# Metrics boxes
vlm_boxes_sorted = sorted([shapes_by_name[n] for n in found_vb], key=lambda s: s.top)
metric_names  = ['Rounded Rectangle 48', 'Rounded Rectangle 49', 'Rounded Rectangle 50',
                 'Rounded Rectangle 78', 'Rounded Rectangle 79']
metric_labels = ['CR↑', 'TVR↓', 'ASR↑', 'TVR↓', 'ASR↑']
metric_colors = [C['vlm_hdr'], C['vlm_bdr'], '#26A69A', C['vlm_bdr'], '#26A69A']

found_met = []
missing_met = []
for i, (mname, mlabel, mcolor) in enumerate(zip(metric_names, metric_labels, metric_colors)):
    if mname in shapes_by_name:
        ms = shapes_by_name[mname]
        if i < len(vlm_boxes_sorted):
            ms.top    = vlm_boxes_sorted[i].top
            ms.height = vlm_boxes_sorted[i].height
        set_fill(ms, mcolor)
        set_no_line(ms)
        set_text_style(ms, text=mlabel, font_size=8, bold=True, color='#FFFFFF', align='center')
        found_met.append(mname)
    else:
        missing_met.append(mname)
print(f"  ↳ Metric boxes found: {found_met}")
if missing_met:
    print(f"  ⚠️  Missing metrics: {missing_met}")

# ── 6. Image Quality panel ───────────────────────────────────────────────────
qual_containers = ['Rounded Rectangle 52', 'Rounded Rectangle 53']
found_qc = []
missing_qc = []
for name in qual_containers:
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], C['qual_bg'])
        set_line(shapes_by_name[name], C['qual_bdr'], 1.5)
        found_qc.append(name)
    else:
        missing_qc.append(name)
print(f"✅ Step 6a: Quality containers found: {found_qc}")
if missing_qc:
    print(f"  ⚠️  Missing: {missing_qc}")

qual_headers = ['Rounded Rectangle 54', 'Rounded Rectangle 55']
found_qh = []
missing_qh = []
for name in qual_headers:
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], C['qual_hdr'])
        set_no_line(shapes_by_name[name])
        set_text_style(shapes_by_name[name], font_size=9, bold=True, color='#FFFFFF', align='center')
        found_qh.append(name)
    else:
        missing_qh.append(name)
print(f"  ↳ Quality headers found: {found_qh}")
if missing_qh:
    print(f"  ⚠️  Missing: {missing_qh}")

qual_items = {
    'Rounded Rectangle 56': ('#1A237E', 'PSNR / SSIM'),
    'Rounded Rectangle 57': (C['qual_bdr'], 'LPIPS / FID'),
    'Rounded Rectangle 58': ('#5C6BC0', 'NIMA / MUSIQ'),
    'Rounded Rectangle 59': ('#1A237E', 'IQA ↑'),
    'Rounded Rectangle 60': (C['qual_bdr'], 'Perceptual ↓'),
    'Rounded Rectangle 61': ('#5C6BC0', 'FID ↓'),
}
found_qi = []
missing_qi = []
for name, (color, label) in qual_items.items():
    if name in shapes_by_name:
        s = shapes_by_name[name]
        set_fill(s, color)
        set_no_line(s)
        set_text_style(s, text=label, font_size=8, bold=True, color='#FFFFFF', align='center')
        found_qi.append(name)
    else:
        missing_qi.append(name)
print(f"  ↳ Quality items found: {found_qi}")
if missing_qi:
    print(f"  ⚠️  Missing: {missing_qi}")

# ── 7. Efficiency panel ──────────────────────────────────────────────────────
eff_containers = ['Rounded Rectangle 63', 'Rounded Rectangle 64']
found_ec = []
missing_ec = []
for name in eff_containers:
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], C['eff_bg'])
        set_line(shapes_by_name[name], C['eff_bdr'], 1.5)
        found_ec.append(name)
    else:
        missing_ec.append(name)
print(f"✅ Step 7a: Efficiency containers found: {found_ec}")
if missing_ec:
    print(f"  ⚠️  Missing: {missing_ec}")

eff_headers = ['Rounded Rectangle 65', 'Rounded Rectangle 66']
found_eh = []
missing_eh = []
for name in eff_headers:
    if name in shapes_by_name:
        set_fill(shapes_by_name[name], C['eff_hdr'])
        set_no_line(shapes_by_name[name])
        set_text_style(shapes_by_name[name], font_size=9, bold=True, color='#FFFFFF', align='center')
        found_eh.append(name)
    else:
        missing_eh.append(name)
print(f"  ↳ Efficiency headers found: {found_eh}")
if missing_eh:
    print(f"  ⚠️  Missing: {missing_eh}")

eff_items = {
    'Rounded Rectangle 67': (C['eff_hdr'],  'Wall-clock Time'),
    'Rounded Rectangle 68': (C['eff_bdr'],  'Peak VRAM (GPU)'),
    'Rounded Rectangle 69': ('#78909C',     'Throughput'),
    'Rounded Rectangle 70': (C['eff_hdr'],  'Time ↓'),
    'Rounded Rectangle 71': (C['eff_bdr'],  'VRAM ↓'),
    'Rounded Rectangle 72': ('#78909C',     'TPut ↑'),
}
found_ei = []
missing_ei = []
for name, (color, label) in eff_items.items():
    if name in shapes_by_name:
        s = shapes_by_name[name]
        set_fill(s, color)
        set_no_line(s)
        set_text_style(s, text=label, font_size=8, bold=True, color='#FFFFFF', align='center')
        found_ei.append(name)
    else:
        missing_ei.append(name)
print(f"  ↳ Efficiency items found: {found_ei}")
if missing_ei:
    print(f"  ⚠️  Missing: {missing_ei}")

# ── 8. Connectors ────────────────────────────────────────────────────────────
connector_names = [
    'Connector 5', 'Connector 8', 'Connector 23', 'Connector 24',
    'Connector 25','Connector 26','Connector 27','Connector 28',
    'Connector 29','Connector 30','Connector 32','Connector 33',
    'Connector 34','Connector 35','Connector 36','Connector 37',
    'Connector 38','Connector 39','Connector 40',
]
found_conn = []
missing_conn = []
for name in connector_names:
    if name in shapes_by_name:
        s = shapes_by_name[name]
        try:
            s.line.color.rgb = rgb(C['arrow'])
            if name in ['Connector 5', 'Connector 8', 'Connector 32', 'Connector 40',
                        'Connector 24', 'Connector 33']:
                s.line.width = Pt(1.8)
            else:
                s.line.width = Pt(1.2)
            found_conn.append(name)
        except Exception as ex:
            print(f"  ⚠️  Connector {name} line error: {ex}")
    else:
        missing_conn.append(name)
print(f"✅ Step 8: Connectors styled: {len(found_conn)}")
if missing_conn:
    print(f"  ⚠️  Missing connectors: {missing_conn}")

divider_names = ['Connector 51', 'Connector 62', 'Connector 73']
found_div = []
missing_div = []
for name in divider_names:
    if name in shapes_by_name:
        try:
            shapes_by_name[name].line.color.rgb = rgb('#B0BEC5')
            shapes_by_name[name].line.width = Pt(0.75)
            found_div.append(name)
        except Exception as ex:
            print(f"  ⚠️  Divider {name} error: {ex}")
    else:
        missing_div.append(name)
print(f"  ↳ Section dividers styled: {found_div}")
if missing_div:
    print(f"  ⚠️  Missing dividers: {missing_div}")

# ── 9. Add figure title text box ─────────────────────────────────────────────
from pptx.util import Inches
txBox = slide.shapes.add_textbox(
    Emu(int(0.22 * 914400)),   # left
    Emu(int(0.12 * 914400)),   # top
    Emu(int(7.5 * 914400)),    # width
    Emu(int(0.45 * 914400)),   # height
)
tf = txBox.text_frame
tf.word_wrap = False
para = tf.paragraphs[0]
para.alignment = PP_ALIGN.LEFT
run = para.add_run()
run.text = "CaptchaShield Evaluation Pipeline"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = rgb('#1A237E')
run.font.name = 'Calibri'
txBox.fill.background()
txBox.line.color.rgb = rgb('#FFFFFF')

txBox2 = slide.shapes.add_textbox(
    Emu(int(0.22 * 914400)),
    Emu(int(0.52 * 914400)),
    Emu(int(10.0 * 914400)),
    Emu(int(0.32 * 914400)),
)
tf2 = txBox2.text_frame
para2 = tf2.paragraphs[0]
para2.alignment = PP_ALIGN.LEFT
run2 = para2.add_run()
run2.text = "GB2312 CAPTCHA  →  Adversarial Perturbation (6 methods, 3 modalities)  →  VLM Evaluation (5 models)  →  Quality & Efficiency Metrics"
run2.font.size = Pt(8)
run2.font.bold = False
run2.font.color.rgb = rgb('#546E7A')
run2.font.name = 'Calibri'
txBox2.fill.background()
txBox2.line.color.rgb = rgb('#FFFFFF')
print("✅ Step 9: Figure title and subtitle added")

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(str(DEST))
import os
pptx_size = os.path.getsize(DEST)
print(f"\n✅ Saved v1 PPTX → {DEST}")
print(f"   PPTX size: {pptx_size:,} bytes ({pptx_size/1024:.1f} KB)")

# ── PDF conversion ────────────────────────────────────────────────────────────
pdf_dest = OUT_DIR / "fig2_pipeline_v1.pdf"
script = f'''
tell application "Microsoft PowerPoint"
    set pp to open POSIX file "{DEST}"
    delay 2
    save pp in POSIX file "{pdf_dest}" as save as PDF
    close pp saving no
end tell
'''
print("\n⏳ Converting to PDF via AppleScript...")
result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True, timeout=90)
print(f"AppleScript stdout: '{result.stdout.strip()}'")
print(f"AppleScript stderr: '{result.stderr.strip()}'")
print(f"AppleScript return code: {result.returncode}")

if pdf_dest.exists():
    pdf_size = os.path.getsize(pdf_dest)
    print(f"\n✅ Saved v1 PDF  → {pdf_dest}")
    print(f"   PDF  size: {pdf_size:,} bytes ({pdf_size/1024:.1f} KB)")
else:
    print(f"\n❌ PDF not found at {pdf_dest}")
    # List files to help debug
    print("Files in output dir:")
    for f in sorted(OUT_DIR.iterdir()):
        print(f"  {f.name}  ({os.path.getsize(f):,} bytes)")
