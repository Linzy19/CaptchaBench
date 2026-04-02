"""
figmm_fig2_pipeline_pptx.py — ACM MM 2026 Figure 2: Pipeline
=============================================================
Reconstructs the pipeline figure based on "fig2_pipeline.pptx" template
(located in figures/).  Updates content for ACM MM 2026:
  - Dataset:    GB2312 Level-1 / 3,500 Characters → 840,000 images total
                (2 generators × 3,500 × 120)
  - Generation: Stable Diffusion (ID ControlNet) + SDXL ControlNet
  - CAPTCHA:    840,000 images (3,500 × 2 × 120)
  - Perturbation: same 6 methods / 3 modality groups (unchanged)
  - Evaluation: 5 VLMs (adds Qwen-VL + Kimi 2.5 before existing 3),
                same metric groups

Template shape inventory (confirmed by inspection):
  RR 3  → Dataset content box          "GB2312 Level-1\\n3,500 Characters"
  RR 6  → Generation content box       "Illusion Diffusion Pipeline (120 BG)"
  RR 9  → CAPTCHA content box          "CAPTCHA Images\\n420,000\\n(3,500×120)"
  RR 45 → VLM row-1 box                "Gemini 3.0"
  RR 46 → VLM row-2 box                "GPT-5.2"
  RR 47 → VLM row-3 box                "GLM-4V"
  RR 48 → Metric row-1 box             "CR↑"
  RR 49 → Metric row-2 box             "TVR↓"
  RR 50 → Metric row-3 box             "ASR↑"
  RR 41 → VLM Evaluation left container
  RR 42 → VLM Evaluation right (Metrics) container
"""

import copy
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ── Paths ─────────────────────────────────────────────────────────────────────
SRC_DIR       = Path("/Users/lzy/Code_lzy/Posion_paper/ACM_MM26/main/figures")
OUT_DIR       = Path("/Users/lzy/Code_lzy/Posion_paper/ACM_MM26/main/figures/cc_fig")
OUT_DIR.mkdir(parents=True, exist_ok=True)

TEMPLATE_PATH = SRC_DIR / "fig2_pipeline.pptx"
PPTX_PATH     = OUT_DIR / "fig2_pipeline.pptx"
PDF_PATH      = OUT_DIR / "fig2_pipeline.pdf"

# ── Load template ─────────────────────────────────────────────────────────────
prs   = Presentation(str(TEMPLATE_PATH))
slide = prs.slides[0]
shapes_by_name = {s.name: s for s in slide.shapes}

print(f"Loaded template: {TEMPLATE_PATH.name}")
print(f"  Slide size : {prs.slide_width/914400:.3f}\" × {prs.slide_height/914400:.3f}\"")
print(f"  Shapes     : {len(slide.shapes)}")

# ─────────────────────────────────────────────────────────────────────────────
# Utility: replace text lines in a shape, preserving XML formatting
# ─────────────────────────────────────────────────────────────────────────────

def replace_text_lines(shape_name, lines, bold=None, font_size_pt=None, color_hex=None):
    """
    Replace the text content of *shape_name* with *lines* (list[str]).
    Each element → one <a:p>.  Paragraph/run formatting is cloned from the
    first original paragraph; optional overrides are applied.
    """
    if shape_name not in shapes_by_name:
        print(f"  ⚠  '{shape_name}' not found – skipped")
        return
    shape = shapes_by_name[shape_name]
    if not hasattr(shape, "text_frame"):
        print(f"  ⚠  '{shape_name}' has no text_frame – skipped")
        return

    tf     = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody
    paras  = txBody.findall(qn('a:p'))

    # Snapshot a template run for cloning
    template_r_xml = None
    for p in paras:
        for r in p.findall(qn('a:r')):
            template_r_xml = copy.deepcopy(r)
            break
        if template_r_xml is not None:
            break

    def _make_para(text):
        new_p = copy.deepcopy(paras[0])
        for child in new_p.findall(qn('a:r')):
            new_p.remove(child)
        for child in new_p.findall(qn('a:br')):
            new_p.remove(child)

        r_el = copy.deepcopy(template_r_xml) if template_r_xml is not None else None
        if r_el is None:
            from lxml import etree
            r_el = etree.Element(qn('a:r'))

        t_el = r_el.find(qn('a:t'))
        if t_el is None:
            from lxml import etree
            t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = text

        rPr = r_el.find(qn('a:rPr'))
        if rPr is None:
            from lxml import etree
            rPr = etree.SubElement(r_el, qn('a:rPr'))
            r_el.insert(0, rPr)
        if bold is not None:
            rPr.set('b', '1' if bold else '0')
        if font_size_pt is not None:
            rPr.set('sz', str(int(font_size_pt * 100)))
        if color_hex is not None:
            for sf in rPr.findall(qn('a:solidFill')):
                rPr.remove(sf)
            from lxml import etree
            sf  = etree.SubElement(rPr, qn('a:solidFill'))
            rgb = etree.SubElement(sf,  qn('a:srgbClr'))
            rgb.set('val', color_hex.lstrip('#'))

        new_p.append(r_el)
        return new_p

    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    for line in lines:
        txBody.append(_make_para(line))

    print(f"  ✓  '{shape_name}' → {lines!r}")


def duplicate_shape(slide, source_shape):
    """Deep-copy *source_shape* and append to slide; return new shape."""
    sp_copy = copy.deepcopy(source_shape._element)
    slide.shapes._spTree.append(sp_copy)
    return slide.shapes[-1]


def set_shape_single_text(shape, text):
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = text
            return
    if tf.paragraphs:
        tf.paragraphs[0].text = text


def set_shape_fill(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bytes.fromhex(hex_color.lstrip('#')))


def fix_duplicate_ids(slide):
    """Assign unique IDs and names to any shapes that share an ID."""
    all_shapes = list(slide.shapes)
    max_id     = max(s.shape_id for s in all_shapes)
    seen_ids   = set()
    fixed      = 0
    for i, s in enumerate(all_shapes):
        sid = s.shape_id
        if sid in seen_ids:
            max_id += 1
            new_name = f"Rounded Rectangle {max_id}"
            cNvPr = s._element.find('.//' + qn('p:cNvPr'))
            if cNvPr is not None:
                cNvPr.set('id',   str(max_id))
                cNvPr.set('name', new_name)
                fixed += 1
        seen_ids.add(sid if sid not in seen_ids else max_id)
    if fixed:
        print(f"  ✓  Fixed {fixed} duplicate shape ID(s)")
    return fixed


# ─────────────────────────────────────────────────────────────────────────────
# ① Dataset — update to 840 K
# ─────────────────────────────────────────────────────────────────────────────
print("\n── ① Dataset ────────────────────────────────────────────────────────────")
replace_text_lines(
    "Rounded Rectangle 3",
    ["GB2312 Level-1", "3,500 Characters", "→ 840,000 Images", "(2 Gen × 120 BG)"],
)

# ─────────────────────────────────────────────────────────────────────────────
# ② Generation — two generators
# ─────────────────────────────────────────────────────────────────────────────
print("\n── ② Generation ─────────────────────────────────────────────────────────")
replace_text_lines(
    "Rounded Rectangle 6",
    ["SD (ID ControlNet)", "+", "SDXL ControlNet", "(120 BG each)"],
)

# ─────────────────────────────────────────────────────────────────────────────
# ③ CAPTCHA — 840 K total
# ─────────────────────────────────────────────────────────────────────────────
print("\n── ③ CAPTCHA ─────────────────────────────────────────────────────────────")
replace_text_lines(
    "Rounded Rectangle 9",
    ["CAPTCHA Images", "840,000", "(3,500 × 2 × 120)"],
)

# ─────────────────────────────────────────────────────────────────────────────
# ⑤ VLM Evaluation — rename 3 existing rows, add 2 new rows
#   Before : Gemini 3.0 | GPT-5.2 | GLM-4V
#   After  : Qwen-VL | Kimi 2.5 | GPT-5.2 | Gemini-3.0 | GLM-4V
# ─────────────────────────────────────────────────────────────────────────────
print("\n── ⑤ VLM Evaluation ────────────────────────────────────────────────────")

replace_text_lines("Rounded Rectangle 45", ["Qwen-VL"])
replace_text_lines("Rounded Rectangle 46", ["Kimi 2.5"])
replace_text_lines("Rounded Rectangle 47", ["GPT-5.2"])

vlm1 = shapes_by_name["Rounded Rectangle 45"]
vlm2 = shapes_by_name["Rounded Rectangle 46"]
vlm3 = shapes_by_name["Rounded Rectangle 47"]
met1 = shapes_by_name["Rounded Rectangle 48"]
met2 = shapes_by_name["Rounded Rectangle 49"]
met3 = shapes_by_name["Rounded Rectangle 50"]
vlm_container = shapes_by_name["Rounded Rectangle 41"]
met_container = shapes_by_name["Rounded Rectangle 42"]

row_pitch     = vlm2.top - vlm1.top
met_row_pitch = met2.top - met1.top
print(f"  row_pitch = {row_pitch} EMU  ({row_pitch/914400:.4f}\")")

# Row 4: Gemini-3.0
g_box = duplicate_shape(slide, vlm3)
g_box.left   = vlm1.left
g_box.top    = vlm1.top + 3 * row_pitch
g_box.width  = vlm1.width
g_box.height = vlm1.height
set_shape_single_text(g_box, "Gemini-3.0")
try:
    set_shape_fill(g_box, str(vlm2.fill.fore_color.rgb))
except Exception:
    pass
print("  ✓  Added 'Gemini-3.0' (row 4)")

# Row 5: GLM-4V
glm_box = duplicate_shape(slide, vlm3)
glm_box.left   = vlm1.left
glm_box.top    = vlm1.top + 4 * row_pitch
glm_box.width  = vlm1.width
glm_box.height = vlm1.height
set_shape_single_text(glm_box, "GLM-4V")
try:
    set_shape_fill(glm_box, str(vlm3.fill.fore_color.rgb))
except Exception:
    pass
print("  ✓  Added 'GLM-4V' (row 5)")

# Metric rows 4 & 5
for row, label in [(4, "CR↑"), (5, "CR↑")]:
    m = duplicate_shape(slide, met3)
    m.left   = met1.left
    m.top    = met1.top + row * met_row_pitch
    m.width  = met1.width
    m.height = met1.height
    set_shape_single_text(m, label)
    print(f"  ✓  Added metric '{label}' (row {row})")

# Expand containers
extra = 2 * row_pitch
vlm_container.height += extra
met_container.height += extra
print(f"  ✓  Expanded containers by {extra} EMU ({extra/914400:.4f}\")")

# ── Fix duplicate shape IDs introduced by deep-copy ──────────────────────────
print("\n── Fix duplicate IDs ────────────────────────────────────────────────────")
fix_duplicate_ids(slide)

# ─────────────────────────────────────────────────────────────────────────────
# Save PPTX
# ─────────────────────────────────────────────────────────────────────────────
prs.save(str(PPTX_PATH))
print(f"\n✅  PPTX saved → {PPTX_PATH}  ({PPTX_PATH.stat().st_size:,} bytes)")

# ─────────────────────────────────────────────────────────────────────────────
# Sanity check
# ─────────────────────────────────────────────────────────────────────────────
print("\n── Sanity check ─────────────────────────────────────────────────────────")
prs2   = Presentation(str(PPTX_PATH))
slide2 = prs2.slides[0]
s2     = {sh.name: sh for sh in slide2.shapes}

checks = [
    ("Rounded Rectangle 3",  "840,000"),
    ("Rounded Rectangle 6",  "SDXL"),
    ("Rounded Rectangle 9",  "840,000"),
    ("Rounded Rectangle 45", "Qwen-VL"),
    ("Rounded Rectangle 46", "Kimi 2.5"),
    ("Rounded Rectangle 47", "GPT-5.2"),
]
all_ok = True
for name, kw in checks:
    sh = s2.get(name)
    txt = sh.text if sh else "(missing)"
    ok  = kw in txt
    all_ok = all_ok and ok
    print(f"  {'✅' if ok else '❌'}  {name:<35} → {txt!r:50}  [{kw}]")

# Check newly added boxes (by text content)
new_found = [(sh.name, sh.text) for sh in slide2.shapes
             if hasattr(sh, 'text') and sh.text in ("Gemini-3.0", "GLM-4V")]
print(f"\n  New VLM boxes : {new_found}")
print(f"  Total shapes  : {len(list(slide2.shapes))}  (expected 76)")

ids   = [sh.shape_id for sh in slide2.shapes]
names = [sh.name for sh in slide2.shapes]
print(f"  Duplicate IDs  : {[x for x in ids   if ids.count(x)   > 1]}")
print(f"  Duplicate names: {[x for x in names if names.count(x) > 1]}")
print(f"\n{'✅ All checks passed' if all_ok else '❌ Some checks FAILED'}")

# ─────────────────────────────────────────────────────────────────────────────
# PDF conversion
# ─────────────────────────────────────────────────────────────────────────────
print("\n── PDF conversion ──────────────────────────────────────────────────────")

def try_lo(binary):
    try:
        r = subprocess.run(
            [binary, "--headless", "--convert-to", "pdf",
             "--outdir", str(OUT_DIR), str(PPTX_PATH)],
            capture_output=True, text=True, timeout=120,
        )
        return r.returncode, r.stdout, r.stderr
    except FileNotFoundError:
        return -1, "", f"'{binary}' not found in PATH"
    except subprocess.TimeoutExpired:
        return -2, "", f"'{binary}' timed out"

pdf_ok = False
for bin_name in ("libreoffice", "soffice"):
    rc, out, err = try_lo(bin_name)
    print(f"[{bin_name}]  returncode={rc}")
    if out.strip():  print("  stdout:", out.strip()[:300])
    if err.strip():  print("  stderr:", err.strip()[:300])
    if rc == 0:
        pdf_ok = True
        break

if not pdf_ok:
    # Try AppleScript / Microsoft PowerPoint
    print("\n  LibreOffice unavailable — trying Microsoft PowerPoint via AppleScript …")
    script = f"""
set outFile to (POSIX file "{PDF_PATH}")
tell application "Microsoft PowerPoint"
    set theDoc to open (POSIX file "{PPTX_PATH}") as alias
    delay 4
    save theDoc in outFile as save as PDF
    delay 3
    close theDoc saving no
end tell
return "done"
"""
    result = subprocess.run(["osascript", "-e", script],
                            capture_output=True, text=True, timeout=60)
    print(f"  AppleScript returncode={result.returncode}")
    if result.stdout.strip(): print("  stdout:", result.stdout.strip())
    if result.stderr.strip(): print("  stderr:", result.stderr.strip()[:300])

    if not PDF_PATH.exists():
        # Fallback: use presentation already open in PowerPoint
        script2 = f"""
set outFile to (POSIX file "{PDF_PATH}")
tell application "Microsoft PowerPoint"
    set theDoc to presentation 1
    save theDoc in outFile as save as PDF
    delay 3
    close theDoc saving no
end tell
return "done"
"""
        result2 = subprocess.run(["osascript", "-e", script2],
                                 capture_output=True, text=True, timeout=60)
        print(f"  Fallback AppleScript returncode={result2.returncode}")
        if result2.stdout.strip(): print("  stdout:", result2.stdout.strip())
        if result2.stderr.strip(): print("  stderr:", result2.stderr.strip()[:300])

if PDF_PATH.exists():
    print(f"\n✅  PDF saved → {PDF_PATH}  ({PDF_PATH.stat().st_size:,} bytes)")
else:
    pdfs = sorted(OUT_DIR.glob("*.pdf"))
    if pdfs:
        print(f"\n⚠   PDF saved as: {[p.name for p in pdfs]}")
    else:
        print("\n⚠   PDF conversion failed — no PDF found.")
        print("    Manual command:")
        print(f"      soffice --headless --convert-to pdf --outdir '{OUT_DIR}' '{PPTX_PATH}'")

# ─────────────────────────────────────────────────────────────────────────────
# Directory listing
# ─────────────────────────────────────────────────────────────────────────────
print(f"\n── Files in {OUT_DIR} ──────────────────────────────────────────────────")
for f in sorted(OUT_DIR.iterdir()):
    if not f.name.startswith('~'):   # skip lock files
        print(f"  {f.name:<44}  {f.stat().st_size:>12,} bytes")
